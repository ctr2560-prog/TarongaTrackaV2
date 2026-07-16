#!/usr/bin/env python3
"""Generate 32 Taronga Tracka pre/post lesson PowerPoint decks."""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

OUT_DIR = os.path.join(os.path.dirname(__file__), '..', 'public', 'resources', 'pptx')
IMG_DIR = os.path.join(os.path.dirname(__file__), '..', 'public', 'images')

os.makedirs(OUT_DIR, exist_ok=True)

# Brand colours
C_DEEP   = RGBColor(0x07, 0x1E, 0x14)
C_MID    = RGBColor(0x1A, 0x52, 0x38)
C_EUCAL  = RGBColor(0x2E, 0x7D, 0x55)
C_CANVAS = RGBColor(0xED, 0xEA, 0xE3)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_GOLD   = RGBColor(0xD9, 0x77, 0x06)
C_MUTED  = RGBColor(0xA8, 0xC4, 0xB2)

SUBJ_COLORS = {
    'science': RGBColor(0x1A, 0x52, 0x38),
    'maths':   RGBColor(0x03, 0x69, 0xA1),
    'english': RGBColor(0x7C, 0x3A, 0xED),
    'pdhpe':   RGBColor(0xBE, 0x18, 0x5D),
}

SUBJ_LABELS = {
    'science': 'Science & Technology',
    'maths':   'Mathematics',
    'english': 'English',
    'pdhpe':   'PDHPE',
}

ANIMALS = ['Tiger', 'Lion', 'Giraffe', 'Gorilla', 'Chimpanzee', 'Dingo', 'Koala', 'Sea Lion']
BRAIN_BREAKS = [
    ('Animal Charades', 'Split into pairs. One person acts out an animal from Taronga Zoo — no sounds! Can your partner guess in 30 seconds?'),
    ('Safari Sketch', 'You have 60 seconds to draw an animal from today\'s visit. Show your class and see if they can guess it!'),
    ('Conservation Countdown', 'Stand up! Call out one fact about any Taronga animal. Sit down once you\'ve shared. Keep going until everyone is seated.'),
    ('Animal Alphabet', 'Starting from A, take turns naming an animal. Skip a letter and you\'re out! How far can the class get?'),
    ('Keeper Questions', 'Imagine you\'re a zookeeper. What are three questions you\'d ask to learn more about an animal\'s health?'),
    ('Habitat Hustle', 'On the count of 3, everyone move like you\'re in a jungle / ocean / savanna / rainforest — teacher calls the habitat!'),
    ('Species Showdown', 'Two animals enter, one wins. Vote: which animal has the BEST adaptation for survival — and defend your answer!'),
    ('Zoo Emoji', 'Using only emojis, describe an animal you saw today. Show a partner — can they decode your message?'),
]

# ─── Content database ───────────────────────────────────────────────────────

CONTENT = {
    'science': {
        2: {
            'topic': 'Living Things & Habitats',
            'outcomes': ['ST2-4LW-S', 'ST2-2DP-T'],
            'prior': ['Animals need food, water, shelter and air', 'Different animals live in different habitats', 'Animals have features that help them survive'],
            'tracka_focus': 'Observing animal features and behaviours; recording what animals eat and how they move',
            'pre_content': [
                ('What do animals need to survive?', ['Think about a pet or an animal you know', 'What does it eat? Where does it live?', 'What would happen if its habitat changed?', 'Share with a partner — 30 seconds each']),
                ('Habitats around the world', ['A habitat is the place where an animal lives', 'Taronga Zoo has animals from rainforests, savannas, and coasts', 'Each animal is adapted to its habitat', 'Discussion: which habitat do you think is hardest to survive in?']),
                ('Animal features & adaptations', ['Strong legs for running away from predators', 'Camouflage to hide from danger', 'Sharp teeth or claws for hunting', 'Look at these animals — what feature stands out?']),
                ('Planning your observations', ['At Taronga, you will observe real animals using Tracka', 'Record: What does the animal DO? What does it look like?', 'Think about: Is this behaviour about finding food, staying safe, or resting?', 'The more detail you write, the more points you earn!']),
            ],
            'li': 'Identify the features and behaviours of animals that help them survive in their habitat.',
            'sc': ['I can name at least two features of an animal and explain how they help it survive', 'I can observe an animal at the zoo and record what I notice', 'I can link an animal\'s behaviour to its needs (food, safety, shelter)'],
            'post_content': [
                ('What did we discover?', ['Share: What was the most surprising thing you saw at the zoo?', 'Did any animal behave in an unexpected way?', 'Which animal scored you the most points — and why?', 'Open discussion — no wrong answers!']),
                ('Animal features we observed', ['Look back at your Tracka observations', 'Which features did you notice most? (teeth, feet, colour, size)', 'Did the animal show any survival behaviours while you watched?', 'Share your best observation with the class']),
                ('Linking back to habitats', ['Think about the animal you observed most closely', 'What habitat would it come from in the wild?', 'Would it survive if its habitat was destroyed? Why?', 'This connects to what scientists call conservation']),
            ],
            'action': ('Animal Habitat Diorama', 'Choose one animal you observed at Taronga Zoo. Create a 3D diorama (shoebox, cardboard, natural materials) showing its wild habitat. Include: the animal, 3 things it needs to survive, and a fact card about threats to this habitat. Be ready to present your diorama to the class.'),
            'exit_q': ['Name one feature of the animal you observed and explain WHY it helps the animal survive.', 'Rate your confidence: 🌱 Still learning / 🌿 Getting there / 🌳 Got it!'],
        },
        3: {
            'topic': 'Ecosystems & Biodiversity',
            'outcomes': ['ST3-4LW-S', 'ST3-5LW-T'],
            'prior': ['Food chains show energy flow between organisms', 'Ecosystems include living and non-living things', 'Biodiversity means variety of life in an area'],
            'tracka_focus': 'Recording feeding behaviours, social interactions, and ecosystem roles of multiple species',
            'pre_content': [
                ('Food webs at Taronga', ['Every animal at the zoo has a role in a food web', 'Producers → Primary Consumers → Secondary Consumers', 'What happens if one species disappears?', 'Draw a food chain with 4 animals — trade with a partner and check each other\'s work']),
                ('Biodiversity matters', ['Taronga Zoo protects over 4,000 animals across 350 species', 'Biodiversity = the variety of life in an ecosystem', 'More biodiversity = stronger, more resilient ecosystems', 'Discussion: Why might a zoo help protect wild biodiversity?']),
                ('Ecosystem roles', ['Every organism has a role: producer, consumer, decomposer', 'Predators keep prey populations in balance', 'Keystone species have an outsized impact on their ecosystem', 'Which animals at Taronga might be keystone species?']),
                ('Tracka mission briefing', ['Today you will observe animals across multiple enclosures', 'For each animal: record its likely role in a food web', 'Note any behaviours that tell you about its diet or social structure', 'Try to build a simple food web from your observations']),
            ],
            'li': 'Describe the roles of organisms in an ecosystem and explain how biodiversity supports ecosystem health.',
            'sc': ['I can construct a food web using animals from Taronga Zoo', 'I can explain what would happen if one species was removed', 'I can use my Tracka observations as evidence in my science explanation'],
            'post_content': [
                ('Class ecosystem map', ['Let\'s build a food web together from our Tracka observations', 'Which animals were producers / consumers / predators?', 'Did anyone observe a feeding behaviour? What did you see?', 'Teacher: map this on the whiteboard with class input']),
                ('Biodiversity at risk', ['Many of Taronga\'s animals are endangered in the wild', 'Tiger: ~600 left in the wild. Gorilla: ~1,000 left.', 'What threatens these animals? (habitat loss, poaching, climate)', 'Discuss: what can we do from our school?']),
                ('Evidence from observations', ['Look back at your highest-scoring Tracka observation', 'What made it a strong scientific observation?', 'Peer share: swap observations and give one piece of feedback']),
            ],
            'action': ('Ecosystem Explainer Video', 'In groups of 3–4, create a 60-second explainer video about one animal you observed at Taronga. Your video must include: the animal\'s role in its ecosystem, at least one threat it faces, and one action people can take to help. Upload to your class folder and vote on the most persuasive video.'),
            'exit_q': ['Draw a food chain with 3 organisms that includes an animal you observed at the zoo.', 'Rate your understanding: I can explain ecosystem roles → Yes / Nearly / Not yet'],
        },
        4: {
            'topic': 'Ecosystems, Adaptations & Classification',
            'outcomes': ['SC4-14LW', 'SC4-15LW'],
            'prior': ['Classification groups organisms by shared features', 'Adaptations are inherited traits that improve survival', 'Ecosystems involve interdependent biotic and abiotic factors'],
            'tracka_focus': 'Classifying observed animals, identifying structural and behavioural adaptations, recording ecosystem interactions',
            'pre_content': [
                ('Classification at the zoo', ['Linnaeus\'s system: Kingdom → Phylum → Class → Order → Family → Genus → Species', 'All mammals share: warm-blooded, hair/fur, live young, mammary glands', 'Taronga has mammals, birds, reptiles, fish, and invertebrates', 'Choose an animal you\'ll observe — predict its full classification']),
                ('Structural vs behavioural adaptations', ['Structural: physical features (beak shape, fur thickness, eye position)', 'Behavioural: actions (migration, camouflage behaviour, nocturnal activity)', 'Both are inherited — encoded in DNA, passed to offspring', 'Discussion: give an example of each from an animal you know']),
                ('Biotic and abiotic factors', ['Biotic: living components — prey, predators, parasites, competitors', 'Abiotic: non-living — temperature, rainfall, soil, sunlight', 'An organism\'s niche is its role + requirements in an ecosystem', 'How does a zoo try to replicate abiotic factors from the wild?']),
                ('Tracka observation strategy', ['For each animal: record 3 structural adaptations and 1 behavioural adaptation', 'Link each adaptation to a survival advantage', 'Note: is the animal solitary or social? What does this tell us?', 'Your written observations are marked on detail — use scientific vocabulary']),
            ],
            'li': 'Analyse the structural and behavioural adaptations of animals and link them to survival advantages in their ecosystem.',
            'sc': ['I can identify and classify an animal using biological taxonomy', 'I can describe at least two adaptations and explain their survival value', 'I can write a scientific observation using precise vocabulary'],
            'post_content': [
                ('Adaptation analysis', ['Share: what was the most impressive adaptation you observed?', 'Could this adaptation exist in a different environment? Why/why not?', 'Peer challenge: name an adaptation — class guesses which animal has it']),
                ('Conservation science', ['Taronga\'s conservation work uses science to protect endangered species', 'Breeding programs require understanding genetics and adaptation', 'Discuss: how does loss of genetic diversity threaten a species?', 'Which animals at Taronga are part of breeding programs?']),
                ('Evaluating our observations', ['Review your Tracka scores — what domain was weakest?', 'How could you improve your scientific observation writing?', 'What question do you still have about one of the animals you observed?']),
            ],
            'action': ('Adaptation Field Report', 'Write a 400–500 word scientific report on one animal you observed at Taronga Zoo. Structure: Introduction (classification + habitat), Body (3 adaptations with diagrams), Analysis (how adaptations link to survival), Conclusion (conservation status + threats). Include your Tracka observation as primary evidence.'),
            'exit_q': ['Name one structural and one behavioural adaptation of your chosen animal. Explain the survival advantage of each.', '⭐ Extension: How might climate change affect this adaptation over generations?'],
        },
        5: {
            'topic': 'Evolution, Natural Selection & Conservation',
            'outcomes': ['SC5-14LW', 'SC5-15LW'],
            'prior': ['Natural selection drives evolutionary change over generations', 'Variation within populations is heritable', 'Extinction occurs when a species cannot adapt fast enough'],
            'tracka_focus': 'Evaluating conservation status, analysing adaptation as evidence of evolution, critiquing human impact on wild populations',
            'pre_content': [
                ('Evidence for evolution', ['Fossil record, comparative anatomy, DNA analysis all support evolution', 'Homologous structures: same origin, different function (human arm / whale flipper / bat wing)', 'Vestigial structures: remnants of ancestral features (human tailbone, whale pelvis)', 'Discussion: what evidence could you find at the zoo?']),
                ('Natural selection in action', ['Variation → Selection pressure → Differential survival → Reproduction → Change over generations', 'This is NOT goal-directed — it\'s statistical', 'Example: antibiotic resistance evolves in bacteria within years', 'How might captivity change selection pressures on zoo animals?']),
                ('Conservation genetics', ['Small populations → inbreeding → reduced genetic diversity → reduced fitness', 'Taronga participates in Species Survival Plans (SSP) and studbooks', 'De-extinction debate: should we try to bring back extinct species?', 'What are the ethical implications of captive breeding programs?']),
                ('Critical observation brief', ['Today you will gather primary evidence of evolutionary adaptations', 'For each animal: hypothesise the selection pressure that drove each adaptation', 'Consider: does captivity expose the animal to the same pressures as the wild?', 'This data will form the evidence base for your action project']),
            ],
            'li': 'Evaluate how natural selection drives adaptation and assess the role of conservation science in managing evolutionary fitness of endangered species.',
            'sc': ['I can explain natural selection using variation, selection pressure, and differential reproduction', 'I can analyse an animal\'s adaptations as products of evolutionary history', 'I can construct a reasoned argument about conservation ethics using scientific evidence'],
            'post_content': [
                ('Evolution evidence debrief', ['What adaptations did you observe that most clearly show evolutionary history?', 'Challenge: identify one potential vestigial structure in an animal you saw', 'Discuss: Is a zoo animal "still evolving"? What does the evidence suggest?']),
                ('Conservation science critique', ['Taronga\'s work is scientifically significant — but is it enough?', 'Success stories: Arabian Oryx, Lord Howe Island Stick Insect recovered from captive programs', 'Limitations: zoo animals can lose wild behaviours, genetic diversity is still reduced', 'Class debate: Is captive conservation a solution or a band-aid?']),
                ('Reflection on primary data', ['Compare your Tracka observation score to a peer\'s — what made the difference?', 'Identify one claim in your observation that needs stronger scientific evidence', 'How would a field biologist\'s observation differ from yours?']),
            ],
            'action': ('Conservation Policy Brief', 'Write a 600–800 word policy brief addressed to the NSW Environment Minister arguing for or against expanding Taronga\'s captive breeding program for one critically endangered species you observed. Include: evolutionary biology rationale, current population data, ethical considerations, and a clear recommendation. Format as a professional document with headings, references, and your Tracka observation as Appendix A.'),
            'exit_q': ['Explain why genetic diversity matters for a species\' long-term survival. Use an animal you observed as your example.', '⭐ Extension: Evaluate one limitation of using zoo observations as scientific evidence.'],
        },
    },
    'maths': {
        2: {
            'topic': 'Data, Measurement & Patterns',
            'outcomes': ['MA2-DATA-01', 'MA2-GM-02', 'MA2-RN-01'],
            'prior': ['Collecting and displaying data in tables and graphs', 'Measuring length, mass and time using formal units', 'Identifying and continuing patterns with numbers'],
            'tracka_focus': 'Counting and sorting animals by category, measuring distances on the zoo map, recording data in tally charts',
            'pre_content': [
                ('Collecting data at the zoo', ['Data is information we collect to answer a question', 'We can collect data by counting, measuring, or observing', 'At the zoo: How many animals are in each enclosure? What do they eat?', 'Activity: design a tally chart to count animals by type (mammal, bird, reptile)']),
                ('Graphs and displays', ['Column graphs show totals — each column is a category', 'Picture graphs use symbols — each symbol = one (or more) items', 'We will create a graph of animals we observe at Taronga', 'Look at this data — what does it tell us? What questions could we ask?']),
                ('Measuring at the zoo', ['Length: how far from the viewing platform to the enclosure?', 'Mass: a gorilla weighs about 180 kg. What is that in grams?', 'Time: how long does a sea lion spend in the water vs on land?', 'Challenge: estimate, then check — which was closest?']),
                ('Pattern mission', ['Taronga Tracka awards points in patterns: 10, 20, 40, 80...', 'What type of pattern is this? (doubling)', 'If I score 40 on Tiger and 80 on Lion, what might I score on Gorilla?', 'Make your own animal points pattern and share with a partner']),
            ],
            'li': 'Collect, organise and display data from real-world animal observations and identify number and measurement patterns.',
            'sc': ['I can create a tally chart and column graph from zoo data', 'I can measure and compare lengths and masses using formal units', 'I can identify and continue a number pattern, explaining the rule'],
            'post_content': [
                ('Our class data', ['Let\'s combine all our Tracka observation counts', 'Which animal did the class observe most? Least?', 'Create a class column graph on the whiteboard together', 'What questions could a zoo scientist answer using this data?']),
                ('Measurement discoveries', ['What was the biggest animal you saw? How do we know?', 'Compare: a giraffe is 5.5 metres tall. How many of you stacked up?', 'Distance challenge: using the zoo map, calculate which two enclosures are furthest apart', 'Record your working — show how you found the answer']),
                ('Patterns in animal data', ['Look at your Tracka points — is there a pattern in what scored highest?', 'Can you find a number pattern in animal facts? (teeth, legs, years in captivity)', 'Create a two-step number pattern using an animal fact as your starting number']),
            ],
            'action': ('Animal Data Poster', 'Choose 4 animals from your Taronga visit. Create a data poster showing: a tally chart of their features (eg. number of legs, type of food), a column or picture graph comparing a measurement (height, mass, or lifespan), and one number pattern you found in the data. Present your poster to the class and explain what your data shows.'),
            'exit_q': ['Write one thing you measured or counted at the zoo. Show this as a number sentence.', 'Can you continue this pattern? 5, 10, 20, 40, _____, _____  What\'s the rule?'],
        },
        3: {
            'topic': 'Statistics, Chance & Measurement',
            'outcomes': ['MA3-DATA-01', 'MA3-CHAN-01', 'MA3-GM-02'],
            'prior': ['Displaying data in tables, column graphs and dot plots', 'Describing chance using fractions and percentages', 'Converting units of measurement and calculating perimeter and area'],
            'tracka_focus': 'Analysing class Tracka scores as a data set, calculating chance of achieving specific score ranges, estimating enclosure areas',
            'pre_content': [
                ('Mean, median and mode', ['The mean is the average — add all values, divide by count', 'The median is the middle value when ordered', 'The mode is the most common value', 'Use these with: How long do animals at Taronga live? (class data set)']),
                ('Chance and probability', ['Probability = favourable outcomes ÷ total outcomes', 'If 3 of 10 animals are mammals, chance of seeing a mammal = 3/10 = 30%', 'Express as fraction, decimal or percentage', 'Is it equally likely that any two animals will appear at the same time? Discuss.']),
                ('Measurement and area', ['Taronga Zoo covers 7 hectares (70,000 m²)', 'The tiger enclosure is approx 2,000 m² — calculate its dimensions if rectangular', 'Perimeter vs area — what\'s the difference? When does each matter for zookeepers?', 'Convert: 2 km² = _____ m² = _____ hectares']),
                ('Tracka data strategy', ['You will collect score data at the zoo — use it like a real data set', 'Record: your score per animal AND the time you spent observing', 'After the visit, we will calculate class mean, median and mode for scores', 'Think: is there a relationship between time spent and score?']),
            ],
            'li': 'Collect and interpret statistical data from zoo observations, describe probability using fractions and percentages, and apply measurement concepts to real contexts.',
            'sc': ['I can calculate and interpret mean, median and mode from a data set', 'I can express probability as a fraction, decimal and percentage', 'I can calculate area and perimeter and convert between units'],
            'post_content': [
                ('Class score statistics', ['Enter everyone\'s total Tracka scores on the board', 'Calculate: class mean score, median score, mode score', 'Draw a dot plot of the data — what shape does the distribution have?', 'What does an outlier in our data tell us?']),
                ('Probability at the zoo', ['If Taronga has 350 species, and 40 are mammals — what\'s the probability of visiting a mammal enclosure?', 'In your Tracka quiz: if a question has 4 options, what\'s the chance of guessing correctly?', 'Did you notice any animals doing unexpected things? Was that likely or unlikely?']),
                ('Measurement in context', ['Estimate the area of one enclosure you visited, using steps or visual comparison', 'A keeper prepares 5 kg of food per large mammal per day — how much for a week? A month?', 'Challenge: if the zoo expanded by 20%, what would its new area be?']),
            ],
            'action': ('Zoo Statistics Report', 'You are a junior zookeeper preparing a statistical report for Taronga\'s board. Using your Tracka data and at least 3 other zoo facts, calculate mean/median/mode for two different data sets, display your data in two different graph types, and include a probability question with solution. Present your report as a formal document with a title, sections and conclusion.'),
            'exit_q': ['The class Tracka scores are: 45, 62, 78, 62, 90, 55. Calculate the mean, median and mode.', 'If 8 out of 20 animals at the zoo are from Africa, what is the probability of visiting an African animal first?'],
        },
        4: {
            'topic': 'Data Analysis, Algebra & Financial Maths',
            'outcomes': ['MA4-DAT-C-01', 'MA4-ALG-C-01', 'MA4-FIN-C-01'],
            'prior': ['Constructing and interpreting statistical graphs', 'Using variables and expressions in algebra', 'Calculating percentages, profit/loss and simple interest'],
            'tracka_focus': 'Statistical analysis of observation scores, building algebraic models for point systems, cost analysis of zoo conservation programs',
            'pre_content': [
                ('Statistical displays for data analysis', ['Frequency tables, histograms and box-and-whisker plots', 'Range, mean, median, mode and outliers', 'Taronga context: animal lifespans, population sizes, conservation funding data', 'Examine this data set: Sumatran Tiger wild population by year 2000–2024. What trend do you see?']),
                ('Algebra and animal populations', ['Let p = current population, r = annual growth rate', 'Future population formula: F = p × (1 + r)^n', 'Taronga\'s breeding program adds ~5 tigers per year — model this algebraically', 'Graph your model: x-axis = years, y-axis = population. What does the shape tell us?']),
                ('Financial mathematics in conservation', ['Running Taronga Zoo costs approximately $80 million per year', 'Revenue sources: ticket sales (~$50M), donations, government grants', 'If adult tickets are $45 and child tickets are $25, find the break-even attendance', 'Model: if costs rise 3% per year, what does Taronga need to charge in 5 years?']),
                ('Tracka mathematics mission', ['Your score is calculated: (behaviour + detail + writing) ÷ 15 × 100 + quiz bonus', 'Write this as an algebraic expression using variables b, d, w, q', 'What score do you need to average across 5 animals to reach 400 total points?', 'Plan your strategy — which animals will you prioritise and why?']),
            ],
            'li': 'Apply statistical analysis, algebraic modelling and financial mathematics to real-world conservation and zoo management contexts.',
            'sc': ['I can construct and interpret statistical displays including box plots', 'I can write and evaluate algebraic expressions and formulae', 'I can apply percentage and financial maths to real-world scenarios with justification'],
            'post_content': [
                ('Statistical analysis of class data', ['Enter class Tracka scores — construct a box-and-whisker plot together', 'Identify Q1, Q2 (median), Q3, IQR, and any outliers', 'What does the spread tell us about how consistently the class observed?', 'Compare two subjects\' data — whose distribution is more spread? What does that mean?']),
                ('Modelling with algebra', ['Using the scoring formula: if b = 4, d = 3, w = 4, q = 1 — calculate your score', 'What values of b, d, w maximise your score without the quiz bonus?', 'Challenge: write an inequality showing the conditions under which quiz bonus makes the difference']),
                ('Conservation funding analysis', ['Research task: Taronga\'s WILD LIFE Conservation Fund receives $2.4M annually', 'If 60% goes to field programs and 40% to captive breeding — calculate each allocation', 'If a breeding program needs $180,000 per species — how many species can be funded?', 'Prepare a 2-minute pitch for how you\'d allocate the funds differently']),
            ],
            'action': ('Conservation Finance Model', 'Build a financial model for a new Taronga conservation program. Choose one endangered animal you observed. Research: current wild population, annual program cost, projected outcome over 10 years. Use algebra to model population growth, construct a statistical display of projected data, and calculate the cost per animal saved. Present as a 5-slide deck with your calculations shown.'),
            'exit_q': ['The Tracka formula is: score = (b + d + w) / 15 × 100. If your total score is 73.3, and b = d = 4, find w.', 'A zoo has a budget of $500,000 for 3 programs costing $120k, $200k and $180k. What % of budget is each program?'],
        },
        5: {
            'topic': 'Statistical Analysis, Modelling & Probability',
            'outcomes': ['MA5-DAT-C-01', 'MA5-PRO-C-01', 'MA5-ALG-C-01'],
            'prior': ['Two-way tables, scatter plots and lines of best fit', 'Compound probability and Venn diagrams', 'Non-linear functions and modelling with technology'],
            'tracka_focus': 'Bivariate data analysis of observation scores, modelling exponential population decline, probability of conservation success scenarios',
            'pre_content': [
                ('Bivariate data and correlation', ['Scatter plots show relationships between two variables', 'Correlation: positive, negative, or no correlation', 'Line of best fit (least squares regression) — equation y = mx + b', 'Dataset: time spent observing (x) vs Tracka score (y) — predict the correlation']),
                ('Exponential models in ecology', ['Population models: P(t) = P₀ × e^(rt)', 'If r < 0: population decline (endangered species scenario)', 'Sumatran Tiger: ~600 in 1990, ~400 in 2024. Find the rate r.', 'Project: if decline continues at this rate, when does population reach 100?']),
                ('Compound probability in conservation', ['P(A ∩ B) = P(A) × P(B) for independent events', 'P(A ∪ B) = P(A) + P(B) − P(A ∩ B)', 'Conservation scenario: P(captive breeding succeeds) = 0.7, P(habitat protected) = 0.4', 'Calculate: P(species survives) assuming both are needed. What if either alone is sufficient?']),
                ('Data collection design', ['At Taronga, you will collect bivariate data: time × score, detail × writing', 'Record all variables precisely — we will run a regression analysis after the visit', 'Think about: what extraneous variables might affect your results?', 'Design a simple data collection sheet for your zoo visit']),
            ],
            'li': 'Apply bivariate data analysis, exponential modelling and compound probability to evaluate real-world conservation outcomes with mathematical rigour.',
            'sc': ['I can construct a scatter plot, determine correlation and interpret a regression equation', 'I can build and solve exponential population models', 'I can calculate compound probability and apply it to multi-stage conservation scenarios'],
            'post_content': [
                ('Regression analysis', ['Plot your class data: time observed (x) vs Tracka score (y)', 'Find the line of best fit equation using technology (Desmos / GeoGebra)', 'Interpret the gradient: for each extra minute of observation, score increases by ___', 'Identify outliers — what might explain a point far from the line?']),
                ('Population modelling workshop', ['Using Desmos: plot the Sumatran Tiger population data', 'Fit an exponential model P(t) = a × bᵗ to the data', 'Adjust parameters — which model fits best? (check R² value)', 'Challenge: add a carrying capacity — how does a logistic model differ?']),
                ('Probability decision analysis', ['Taronga must decide between two programs for limited funding', 'Program A: P(success) = 0.8, cost = $300k; Program B: P(success) = 0.5, cost = $150k', 'Calculate expected value for each: E = P(success) × outcome value', 'Which would you fund? Justify using expected value AND ethical reasoning']),
            ],
            'action': ('Mathematical Conservation Report', 'Produce a full mathematical investigation (800–1000 words + working) on one of: (a) Exponential population modelling for an endangered Taronga animal — build, validate, and critique a predictive model; or (b) Statistical analysis of Taronga visitor and conservation funding data — test a hypothesis using a scatter plot, regression, and significance discussion. Include all graphs, equations, and a written interpretation. Submit as a formal report.'),
            'exit_q': ['A population follows P(t) = 600 × 0.94^t. What is the population after 10 years? When will it fall below 200?', 'Two conservation events are independent: P(A) = 0.6, P(B) = 0.7. Find P(A and B) and P(A or B).'],
        },
    },
    'english': {
        2: {
            'topic': 'Observation Writing & Descriptive Language',
            'outcomes': ['EN2-VOCAB-01', 'EN2-COMP-01', 'EN2-URA-01'],
            'prior': ['Using adjectives, verbs and adverbs to add detail to writing', 'Writing in sentences with correct punctuation', 'Identifying features of informational and narrative texts'],
            'tracka_focus': 'Writing detailed animal observations using descriptive language; building vocabulary from keeper cards and enclosure signs',
            'pre_content': [
                ('Wow words for animals', ['Strong verbs: prowl, stalk, groom, lunge, splash, trumpet, slither', 'Precise adjectives: striped, amber-eyed, silky-furred, razor-clawed', 'Adverbs of manner: cautiously, powerfully, silently, playfully', 'Sort these words into a word wall — which animals do they make you think of?']),
                ('What makes a great observation?', ['Describe WHAT you see — specific details, not "it was cool"', 'Use your senses: What do you see? Hear? Notice about its movement?', 'Include: what the animal IS doing + what body part you\'re observing', 'Compare a weak observation vs a strong observation — what\'s different?']),
                ('Sentence structure', ['Simple: The tiger paced along the fence.', 'Compound: The tiger paced the fence, and its amber eyes scanned the crowd.', 'Complex: As the crowd fell silent, the tiger lowered its head and paced slowly.', 'Rewrite this weak sentence using a complex structure: "The gorilla sat."']),
                ('Vocabulary mission prep', ['At the zoo, find 3 words you don\'t know from signs or keeper cards', 'Write what you THINK they mean — then check', 'Use one of your new words in your Tracka observation', 'The more precise your vocabulary, the higher your score!']),
            ],
            'li': 'Use precise vocabulary, vivid verbs and descriptive detail to write engaging animal observations.',
            'sc': ['I can use at least three different types of describing words (adjectives, verbs, adverbs) in my writing', 'I can write a complex sentence about an animal\'s behaviour', 'I can explain the meaning of three new vocabulary words from the zoo'],
            'post_content': [
                ('Sharing our best sentences', ['Read your favourite sentence from your Tracka observation aloud', 'Class: what made it powerful? (verb choice? specific detail? imagery?)', 'Identify the best verb in each shared sentence — create a class verb wall', 'Everyone upgrade one sentence from their observation using a peer\'s suggestion']),
                ('Vocabulary we discovered', ['Share the new words you collected at the zoo', 'Class glossary: write each word, its meaning, and an example sentence', 'Can we group these words by type? (scientific, descriptive, positional)', 'Which word was most useful for describing an animal?']),
                ('Text features of animal writing', ['Where did you see writing at the zoo? (signs, maps, keeper cards)', 'What text features did they use? (headings, bold text, labels, images)', 'How is a zoo sign different from a storybook? Why?', 'How will you use text features in your action project?']),
            ],
            'action': ('Animal Information Report', 'Write an information report (200–300 words) about one animal you observed at Taronga Zoo. Use three sections: Introduction (what animal, where it lives), Features (what it looks like — use 5 precise describing words), and Behaviour (what it does — use 4 strong verbs). Add a labelled diagram. Include at least two vocabulary words you learned at the zoo.'),
            'exit_q': ['Write one sentence about an animal using a strong verb, a precise adjective, AND an adverb.', 'Circle the best word: The lion [walked / prowled / went] through the grass.'],
        },
        3: {
            'topic': 'Informational & Persuasive Writing',
            'outcomes': ['EN3-VOCAB-01', 'EN3-COMP-01', 'EN3-URA-01'],
            'prior': ['Features of informational and persuasive texts', 'Using evidence to support an argument', 'Complex sentences with subordinate clauses'],
            'tracka_focus': 'Writing structured animal observations; crafting persuasive arguments about conservation using zoo visit evidence',
            'pre_content': [
                ('Informational text structures', ['Problem-Solution: This species is endangered because... Taronga is addressing this by...', 'Cause-Effect: Due to habitat loss, gorilla populations have declined by 60%', 'Compare-Contrast: Chimpanzees and gorillas both... however...', 'Identify which structure best suits: a zoo sign / a research report / a letter to the council']),
                ('The language of persuasion', ['Rhetorical question: "Can we really afford to lose another species?"', 'Statistics as evidence: "Over 60% of primate species are now threatened."', 'Expert authority: "According to Taronga\'s chief scientist..."', 'Emotive language: "These magnificent creatures face an uncertain future."']),
                ('Building an argument', ['Claim → Evidence → Explanation → Link back to argument (CEEL structure)', 'Evidence can be: statistics, expert opinion, personal observation, examples', 'At the zoo: your Tracka observation IS evidence — use it', 'Practice: write one CEEL paragraph about why zoos matter']),
                ('Vocabulary for formal writing', ['Replace casual language with formal: "really good" → "highly effective"', 'Hedging language: "This suggests...", "Evidence indicates...", "It appears that..."', 'Connectives for argument: Furthermore, However, In contrast, Consequently', 'Create a list of 10 formal synonyms for words you use too often']),
            ],
            'li': 'Write structured informational and persuasive texts using evidence from the zoo visit, formal vocabulary and clear argument structures.',
            'sc': ['I can write a persuasive paragraph using CEEL structure with zoo evidence', 'I can use at least three persuasive techniques in my writing', 'I can use formal vocabulary and connective language to strengthen my argument'],
            'post_content': [
                ('Evidence from our visit', ['What facts did you discover at the zoo that surprised you?', 'Which animal gave you the strongest evidence for a conservation argument?', 'Peer check: is your evidence SPECIFIC or vague? Improve one piece.', 'How could your Tracka observation score be used as evidence?']),
                ('Evaluating persuasive techniques', ['Read this conservation advertisement — identify the techniques used', 'Which technique was most effective? Least effective? Why?', 'Did you notice Taronga using persuasive language at the zoo? Where?', 'Discuss: can persuasive writing be dishonest? Where is the line?']),
                ('Feedback and redrafting', ['Swap your CEEL paragraph with a partner', 'Feedback: one strength + one specific improvement', 'Redraft your paragraph based on the feedback', 'Read both versions aloud — which is stronger and why?']),
            ],
            'action': ('Conservation Argument Essay', 'Write a 350–500 word persuasive essay arguing for one conservation action that should be taken to protect an animal you observed at Taronga. Use at least three different persuasive techniques, CEEL structure for each body paragraph, evidence from your zoo visit, and formal vocabulary throughout. Include a title, introduction, two body paragraphs, and a conclusion.'),
            'exit_q': ['Write a one-sentence claim about conservation using an animal from the zoo. Then write one piece of evidence to support it.', 'Which persuasive technique is used here? "Can we really let the last 600 tigers disappear on our watch?"'],
        },
        4: {
            'topic': 'Analytical Writing & Close Reading',
            'outcomes': ['EN4-ECA-C-01', 'EN4-ECB-C-01', 'EN4-URA-C-01'],
            'prior': ['Analysing language choices for effect', 'Identifying purpose and audience in texts', 'Writing analytical paragraphs with textual evidence'],
            'tracka_focus': 'Analysing the language of zoo conservation materials; using first-person observation as primary text for analytical writing',
            'pre_content': [
                ('Text purpose and audience', ['Texts are constructed with a specific purpose and audience in mind', 'A zoo sign differs from a scientific report — same content, different language choices', 'Identify: purpose (inform / persuade / entertain / instruct) and audience for each text type at Taronga', 'How does audience shape vocabulary, sentence length, and tone?']),
                ('Language analysis framework (TEEL)', ['Topic sentence: clear analytical claim about the text', 'Evidence: quote or specific reference to the text', 'Explanation: unpack the language choice and its effect', 'Link: connect back to the overall argument', 'Practice: analyse one sentence from a zoo conservation campaign']),
                ('Figurative language for effect', ['Metaphor: "The tiger is a ghost in the grass"', 'Personification: "The forest mourns its lost inhabitants"', 'Alliteration, assonance, imagery — all create emotional effect', 'At the zoo: collect examples of figurative language from signs, videos, keeper cards']),
                ('First-person observation as text', ['Your Tracka observation is a TEXT — written with a purpose and audience', 'Analysing your own writing: what language choices did YOU make?', 'Effective observation writing uses specific nouns, active verbs, sensory detail', 'Today: treat your Tracka observation as primary evidence for an English essay']),
            ],
            'li': 'Analyse how language choices create meaning and effect in conservation texts, and apply analytical writing skills to evaluate zoo observation writing.',
            'sc': ['I can write an analytical TEEL paragraph with a quote and language explanation', 'I can identify and explain the effect of figurative language in a text', 'I can analyse my own observation writing and identify specific language choices'],
            'post_content': [
                ('Close reading of zoo texts', ['Examine this Taronga conservation sign/video transcript as a class', 'Underline: figurative language, specific word choices, structural features', 'TEEL practice: write one paragraph analysing how the text persuades its audience', 'Share and compare — did different students notice different techniques?']),
                ('Your observation as text', ['Reread your Tracka observation as if you were marking someone else\'s work', 'Identify: strong verb choices, precise nouns, hedging language, figurative language', 'What purpose does your observation serve? Who is the audience?', 'Revise one section to make your language choices more deliberate']),
                ('Comparative analysis', ['Compare a children\'s zoo brochure with a scientific field report on the same animal', 'How does purpose and audience change: vocabulary, sentence structure, formality, evidence type', 'Which text is more persuasive? More accurate? More accessible? Can a text be all three?']),
            ],
            'action': ('Language Analysis Essay', 'Write a 500–600 word analytical essay examining how language is used in Taronga Zoo\'s conservation communication. Select two texts (eg. a keeper\'s sign, a campaign video transcript, or your own Tracka observation). Analyse: figurative language, word choice, sentence structure, and how these choices position the audience. Use TEEL paragraphs throughout. Your Tracka observation may be included as a primary text.'),
            'exit_q': ['Identify the language technique and explain its effect: "The last tiger paces in a shrinking world, its amber eyes searching for what no longer exists."', 'Write a TEEL topic sentence for an essay arguing that conservation language is deliberately emotive.'],
        },
        5: {
            'topic': 'Critical Analysis & Multimodal Texts',
            'outcomes': ['EN5-ECA-C-01', 'EN5-ECB-C-01', 'EN5-URA-C-01'],
            'prior': ['Critical analysis of ideology and positioning in texts', 'Evaluating multimodal elements and their interaction', 'Composing extended analytical and argumentative texts'],
            'tracka_focus': 'Critically analysing Taronga\'s conservation narrative as ideological text; creating multimodal persuasive texts using zoo visit evidence',
            'pre_content': [
                ('Critical reading: ideology in texts', ['Every text embeds a worldview — conservation texts are no exception', 'Whose perspective is centred? Whose is absent? (local communities near endangered habitats?)', 'Taronga\'s narrative: Western zoo as conservation hero — is this the only story?', 'Critical question: what does Taronga\'s language reveal about assumptions regarding animals, nature, and human responsibility?']),
                ('Multimodal text analysis', ['Multimodal texts use multiple modes: visual, linguistic, audio, spatial, gestural', 'Composition: where are elements placed? What is foregrounded?', 'Colour, font, image selection — all carry ideological weight', 'Analyse Taronga\'s website homepage: what story does the visual design tell BEFORE you read a word?']),
                ('Constructing a critical argument', ['Critical essays don\'t just describe — they evaluate and challenge', 'Structure: position (thesis) + counterargument + rebuttal + synthesis', 'Using academic language: "This positions the audience to...", "One might argue... however...", "This representation marginalises..."', 'Draft a thesis: Is Taronga\'s conservation narrative empowering or paternalistic?']),
                ('Primary evidence from your zoo visit', ['At Taronga: collect examples of how the zoo frames its conservation narrative', 'Note: language on signs, images used, animal stories told (and untold)', 'Your Tracka observation is a text you produced WITHIN that institutional framework', 'Reflect: were you positioned to observe in a particular way? How?']),
            ],
            'li': 'Critically analyse how ideology is constructed in conservation texts and compose sophisticated multimodal arguments using evidence from first-hand zoo observation.',
            'sc': ['I can identify and analyse ideological positioning in a multimodal conservation text', 'I can construct a critical argument with counterargument and synthesis', 'I can evaluate my own Tracka observation as an ideologically positioned text'],
            'post_content': [
                ('Critical analysis debrief', ['What ideology underpins Taronga\'s conservation narrative? Share your findings.', 'Challenge: is it possible to communicate conservation without ideological positioning?', 'Discuss: who benefits from Taronga\'s narrative, and who might be disadvantaged?', 'Has this changed how you think about zoos and conservation? Why/why not?']),
                ('Multimodal text workshop', ['Share your draft multimodal composition with a peer', 'Peer analysis: what ideology does YOUR text embed? Is this intentional?', 'Evaluate: how effectively do your visual and linguistic choices work TOGETHER?', 'Revise one element specifically to strengthen ideological clarity']),
                ('Reflective writing', ['Write for 5 minutes: How did this excursion change (or reinforce) your assumptions about zoos, animals, and conservation?', 'Identify one assumption you held BEFORE the visit that was challenged', 'How might you write differently about animals after this experience?']),
            ],
            'action': ('Critical Multimodal Essay & Text', 'Produce: (a) A 700–900 word critical essay analysing ideology in TWO Taronga Zoo texts (one visual/digital, one written), examining how they construct conservation as a concept and position their audience — use TEEL paragraphs, critical vocabulary, and your zoo visit as evidence; AND (b) A short original multimodal text (poster, video script, or digital composition) that presents an alternative conservation narrative, with a 150-word reflection on the ideological choices you made.'),
            'exit_q': ['Explain how one visual or language choice in a Taronga text positions its audience ideologically.', '⭐ Extension: Why is it impossible to create a "neutral" conservation text? Argue in 3 sentences.'],
        },
    },
    'pdhpe': {
        2: {
            'topic': 'Health, Movement & Wellbeing',
            'outcomes': ['PD2-2', 'PD2-4', 'PD2-10'],
            'prior': ['Healthy habits include exercise, sleep, food, and water', 'Different movements build different skills', 'Emotions affect our wellbeing and behaviour'],
            'tracka_focus': 'Observing animal movement patterns and linking to health and fitness concepts; connecting animal behaviours to human wellbeing practices',
            'pre_content': [
                ('Animals and movement', ['Different animals move in different ways for different reasons', 'Running, swimming, climbing, slithering — each needs different body adaptations', 'Think about your PE lessons — which animal movement could you copy?', 'Activity: copy these animal movements (5 seconds each): gorilla walk, sea lion flap, tiger stalk']),
                ('Health in the wild', ['Wild animals need food, water, shelter, sleep — just like us', 'Exercise: animals in the wild move all day to hunt, forage, escape predators', 'Stress: animals at the zoo can show signs of stress (pacing, repetitive behaviour)', 'Discuss: how is an animal\'s health similar to your health? How is it different?']),
                ('Emotions and behaviour', ['Animals show emotions through body language', 'Happy / playful / stressed / scared — can you tell which is which?', 'Humans also show emotions in body language — how?', 'At the zoo: look for signs of the animal\'s emotional state in its behaviour']),
                ('Tracka observation focus', ['For PDHPE: focus on HOW the animal moves and WHY', 'Is the movement for exercise? Safety? Social connection? Play?', 'Link what you see to a human health concept', 'The stronger the connection you make, the better your observation score']),
            ],
            'li': 'Identify how animals meet their health needs through movement and behaviour, and connect these to our own health and wellbeing practices.',
            'sc': ['I can describe how an animal moves and what need this movement meets', 'I can identify at least two ways animals and humans have similar health needs', 'I can observe animal body language and describe what emotion it might show'],
            'post_content': [
                ('Movement observations', ['Share: what was the most impressive animal movement you saw?', 'Which movement could you try in PE? Let\'s try it now — 30 seconds each!', 'Compare: gorilla walk vs sea lion movement — what\'s different about how they use their body?', 'Class vote: which animal is the best "athlete"?']),
                ('Health and environment', ['Did the animals seem healthy? What made you think that?', 'What do zookeepers do to keep animals healthy? (vet care, enrichment, diet)', 'Connection: zookeepers are like PE teachers/nurses for animals!', 'What can animals teach us about looking after our own health?']),
                ('Emotional wellbeing', ['Did you see any animals that seemed stressed or bored?', 'What does enrichment mean for zoo animals? (puzzle feeders, new objects, social time)', 'Humans need enrichment too! What\'s YOUR version of enrichment?', 'Discuss: how does school provide enrichment for your brain and body?']),
            ],
            'action': ('Health & Movement Poster', 'Create an A3 poster comparing the health and movement needs of one Taronga Zoo animal to a human. Include: 3 types of movement each animal/human does and why, 2 emotions shown through body language, a "day in the life" comparing the animal\'s routine to yours, and one thing YOU could do to support this animal\'s health (donate, raise awareness, change a habit). Make it colourful and include drawings or printed images.'),
            'exit_q': ['Name one movement you saw an animal doing today. Why do you think it was moving that way?', 'Write one way your health is similar to an animal\'s health and one way it\'s different.'],
        },
        3: {
            'topic': 'Physical Activity, Movement Skills & Health Choices',
            'outcomes': ['PD3-1', 'PD3-4', 'PD3-10'],
            'prior': ['Fundamental movement skills and how they apply to sport/activity', 'The importance of physical activity for physical and mental health', 'Making healthy food and activity choices'],
            'tracka_focus': 'Analysing animal movement for biomechanical principles; connecting wildlife conservation to human health choices and environmental health',
            'pre_content': [
                ('Animal biomechanics', ['Centre of gravity: why do gorillas walk with knuckles down?', 'Force and power: how does a cheetah accelerate to 100 km/h in 3 seconds?', 'Flexibility and range of motion: why can a sea lion rotate its flippers forward?', 'Apply biomechanics language to your favourite sport — what\'s similar?']),
                ('Movement for different purposes', ['Predators move for stealth, speed, and power', 'Prey animals move for agility, vigilance, and stamina', 'Social animals move for communication and bonding (play, grooming)', 'Classify these movements: sprint / patrol / play / forage / escape — predator or prey?']),
                ('Health and conservation', ['Environmental health and human health are connected', 'Air quality, water safety, biodiversity loss all affect human wellbeing', 'Concept: "One Health" — animal health, human health, and ecosystem health are inseparable', 'Discussion: if the tigers go extinct, does it affect your health? How?']),
                ('Observation strategy for PDHPE', ['Focus on: movement quality, apparent fitness, social behaviour, signs of stress', 'Rate the animal on: range of movement / activity level / apparent health (1–5 scale)', 'Note: is the movement voluntary (play, exploration) or stereotypic (pacing, repetitive)', 'Connect each observation to a human health or movement concept']),
            ],
            'li': 'Analyse animal movement using biomechanical principles and explore the links between environmental health, wildlife conservation, and human wellbeing.',
            'sc': ['I can analyse animal movement using biomechanical vocabulary (force, balance, flexibility, power)', 'I can explain the "One Health" concept using a specific example', 'I can design a physical activity inspired by an animal movement I observed'],
            'post_content': [
                ('Biomechanics debrief', ['Which animal movement was biomechanically most impressive? Why?', 'Apply a biomechanical principle (levers, force, balance) to explain it', 'Challenge: design a human exercise drill inspired by this animal movement', 'Lead the class in a 2-minute drill — explain the biomechanical focus']),
                ('One Health discussion', ['Share: how did the zoo change your thinking about the link between nature and health?', 'Research scenario: a local wetland is drained for housing. What happens to human health nearby?', 'Discuss: what health choices can individuals make that also benefit wildlife?', 'Create a class "One Health Pledge" — 3 actions your class will commit to']),
                ('Movement reflection', ['Review your Tracka observation — which aspect of movement did you describe best?', 'Upgrade one sentence using a biomechanical term (leverage, momentum, equilibrium)', 'Share: what was the most surprising movement behaviour you observed?']),
            ],
            'action': ('Animal-Inspired Movement Program', 'Design a 4-week movement/fitness program for Stage 3 students inspired by 4 animals from Taronga Zoo. For each week/animal: describe the animal\'s key movement quality, design 3 fitness activities inspired by this movement (with diagrams), explain the health benefit of this type of movement for humans, and link to the "One Health" concept. Present as a program booklet with a cover page, animal images, clear instructions, and a conclusion about the connection between wildlife and human health.'),
            'exit_q': ['Describe one animal movement you observed today using at least two biomechanical terms.', 'Explain in two sentences how protecting this animal\'s habitat could benefit human health.'],
        },
        4: {
            'topic': 'Movement Concepts, Fitness & Health Advocacy',
            'outcomes': ['PD4-2', 'PD4-5', 'PD4-10'],
            'prior': ['Components of fitness: strength, endurance, flexibility, power, coordination', 'Health promotion models and strategies', 'Analysing the determinants of health'],
            'tracka_focus': 'Applying fitness components to animal analysis; exploring conservation as a form of environmental health advocacy; evaluating determinants of animal and human health',
            'pre_content': [
                ('Components of fitness in wildlife', ['Identify each fitness component in these animals:', 'Tiger: power + agility. Sea lion: cardiovascular endurance + flexibility', 'Gorilla: muscular strength + coordination. Giraffe: balance + flexibility', 'Which component is MOST important for each species\'s survival? Justify.', 'How do humans train for each component? Compare to how animals develop theirs naturally.']),
                ('Determinants of health', ['Biological: genetics, age, sex — tiger\'s immune system, hormones', 'Behavioural: diet, activity, sleep — how does captivity change animal behaviour?', 'Environmental: habitat, pollution, climate — the zoo as controlled environment', 'Socioeconomic: funding, policy — how do government decisions affect wildlife health?']),
                ('Health promotion and conservation', ['The Ottawa Charter for Health Promotion: 5 action areas', 'Map these to conservation: Build Healthy Public Policy (wildlife protection laws), Create Supportive Environments (national parks)', 'Strengthen Community Action (conservation NGOs), Develop Personal Skills (citizen science)', 'Reorient Health Services → Reorient Research (Taronga\'s WILD LIFE Fund)','Discussion: Is conservation a form of public health advocacy?']),
                ('Observation design', ['Focus your Tracka observation on: fitness components evident in movement, health indicators (coat condition, energy level, social engagement), any environmental stressors', 'Use PDHPE vocabulary: evaluate, analyse, justify', 'Connect each observation point to a determinant of health']),
            ],
            'li': 'Analyse components of fitness and determinants of health in animals and evaluate conservation as a form of environmental health advocacy using the Ottawa Charter framework.',
            'sc': ['I can apply at least 4 fitness components to an analysis of animal movement', 'I can explain how biological, behavioural and environmental determinants affect animal health', 'I can map conservation strategies to Ottawa Charter action areas'],
            'post_content': [
                ('Fitness analysis debrief', ['Which animal demonstrated the highest level of physical fitness? Evidence?', 'Did captivity appear to limit any fitness component? How could you tell?', 'Challenge: design one enrichment activity for your chosen animal that targets its weakest fitness component']),
                ('Conservation as health advocacy', ['Map Taronga\'s programs to the Ottawa Charter action areas — which is strongest?', 'What health advocacy skills does a conservation officer need?', 'Discuss: how does environmental degradation become a human health issue? (case study: biodiversity loss → zoonotic disease risk)', 'Class challenge: identify one local environmental action that also improves community health']),
                ('Reflection', ['Review your Tracka observation — how well did it demonstrate PDHPE analytical thinking?', 'What PDHPE vocabulary could strengthen your observation?', 'Has this excursion changed your understanding of health as more than just personal fitness?']),
            ],
            'action': ('Conservation Health Advocacy Campaign', 'Create a health promotion campaign (digital or poster) advocating for the protection of one critically endangered Taronga animal. Apply the Ottawa Charter framework: address all 5 action areas in your campaign. Include: a statistical analysis of the animal\'s population decline, an explanation of how this species\' extinction would affect human health (One Health principle), and a clear call-to-action for your target audience (your school community). Present as a professional campaign package.'),
            'exit_q': ['Name two fitness components visible in an animal you observed. Explain how each aids its survival.', 'Select one Ottawa Charter action area. How could it be applied to protect the animal you studied?'],
        },
        5: {
            'topic': 'Performance Analysis, Health Equity & Advocacy',
            'outcomes': ['PD5-2', 'PD5-5', 'PD5-10'],
            'prior': ['Biomechanical principles applied to performance', 'Health equity and social determinants of health', 'Advocacy skills and community health action'],
            'tracka_focus': 'Critically evaluating captive animal performance and welfare; analysing equity in conservation resource distribution; designing evidence-based advocacy',
            'pre_content': [
                ('Biomechanics and performance analysis', ['Force vectors, torque, Newton\'s laws applied to animal movement', 'Video analysis technique: freeze-frame, identify joint angles, force application', 'Performance analysis in conservation: how do researchers measure animal health?', 'Apply: analyse the biomechanics of a tiger pouncing vs a sea lion swimming — compare force generation and energy efficiency']),
                ('Health equity in conservation', ['Health equity: equal opportunity to achieve full health potential', 'Conservation equity: not all species receive equal funding or attention (flagship vs non-charismatic species)', 'Geographical inequity: most biodiversity in low-income nations; most conservation funding from high-income nations', 'Critical question: who decides which species are worth saving, and by what criteria?']),
                ('Social determinants and "One Health"', ['Social determinants of human health are mirrored in wildlife', 'Poverty → poaching; inequality → land conflict; lack of education → habitat destruction', '"One Health" framework: human, animal, and ecosystem health are inseparable', 'Case study: COVID-19 and zoonotic disease — how does wildlife health directly affect human health?']),
                ('Evidence-based advocacy design', ['Effective advocacy: identifies problem, targets audience, uses evidence, proposes solution, measures impact', 'Taronga as model: examine their advocacy campaigns — what makes them effective?', 'At the zoo: collect evidence (observation, statistics, stories) for your advocacy project', 'Design principle: the most powerful advocacy combines emotion + evidence + clear action']),
            ],
            'li': 'Apply biomechanical analysis to evaluate animal performance, critically examine equity in conservation, and design evidence-based health and conservation advocacy.',
            'sc': ['I can apply Newton\'s laws and biomechanical principles to analyse animal movement performance', 'I can critique equity in conservation funding and decision-making with specific evidence', 'I can design an advocacy campaign that combines evidence, emotional resonance and a clear theory of change'],
            'post_content': [
                ('Performance analysis debrief', ['Share your biomechanical analysis — which animal\'s performance was most complex?', 'How did captivity affect performance capacity? (muscle atrophy, reduced motivation, limited range)', 'Compare: research data on wild vs captive animal fitness outcomes — what do the studies show?', 'Challenge: design a zoo enclosure that would maximise the biomechanical performance of your chosen animal']),
                ('Conservation equity critique', ['Class debate: should conservation funding prioritise endangered mammals or less "charismatic" invertebrates?', 'Who is absent from Taronga\'s conservation narrative? (local communities, Indigenous knowledge holders, non-Western conservationists)', 'Examine: does the zoo experience reinforce or challenge assumptions about animals and nature?', 'Discuss: what would a genuinely equitable global conservation framework look like?']),
                ('Advocacy project development', ['Share your advocacy concept — peer critique using "warm/cool" feedback', 'Evaluate: evidence strength, audience clarity, theory of change, emotional resonance', 'Refine: identify the single most important change to strengthen your advocacy', 'Present final concepts — class votes on most likely to create real change']),
            ],
            'action': ('Evidence-Based Conservation Advocacy Project', 'Produce a complete advocacy project for protecting one critically endangered Taronga animal. Required components: (a) Biomechanical analysis of the animal\'s movement with annotated diagrams; (b) Equity analysis — who are the key stakeholders in this species\' conservation, and whose voices are marginalised?; (c) A "One Health" argument for why protecting this species benefits human health; (d) A full advocacy campaign (platform of your choice: video, social campaign, policy brief, community event plan) with target audience, evidence base, call-to-action, and evaluation framework. Submit with a 200-word critical reflection on your advocacy choices.'),
            'exit_q': ['Apply one of Newton\'s laws to explain a movement you observed at the zoo. Be specific about force, mass and acceleration.', 'Identify one equity issue in wildlife conservation. Who is advantaged? Who is disadvantaged? Explain in 3 sentences.'],
        },
    },
}

def hex_to_rgb(h):
    h = h.lstrip('#')
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def add_rect(slide, l, t, w, h, fill_color, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h, size=18, bold=False, color=C_WHITE, align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = 'Calibri'
    return tb

def add_bullet_text(slide, lines, l, t, w, h, size=16, color=C_WHITE, heading=None, heading_color=None):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    if heading:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = heading
        run.font.size = Pt(size + 2)
        run.font.bold = True
        run.font.color.rgb = heading_color or C_GOLD
        run.font.name = 'Calibri'
    for line in lines:
        p = tf.paragraphs[0] if (first and not heading) else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f'• {line}'
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = 'Calibri'
    return tb

def make_slide_background(slide, prs, color=C_DEEP):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def make_title_slide(prs, subject, stage, timing, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_slide_background(slide, prs, C_DEEP)
    sc = SUBJ_COLORS[subject]

    # Full-width accent bar top
    add_rect(slide, 0, 0, 10, 0.08, sc)

    # Left accent stripe
    add_rect(slide, 0, 0, 0.5, 7.5, sc)

    # Subject colour block top-right
    add_rect(slide, 7.5, 0, 2.5, 2.0, sc)

    # Subject label in block
    add_text(slide, SUBJ_LABELS[subject].upper(), 7.55, 0.1, 2.35, 0.5,
             size=11, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
    add_text(slide, f'STAGE {stage}', 7.55, 0.55, 2.35, 0.4,
             size=22, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
    add_text(slide, timing.upper() + '-VISIT', 7.55, 1.2, 2.35, 0.4,
             size=14, bold=False, color=RGBColor(0xA8,0xC4,0xB2), align=PP_ALIGN.LEFT)

    # Big title
    timing_word = 'Before the Zoo' if timing == 'pre' else 'After the Zoo'
    title = f'{data["topic"]}\n{timing_word}'
    add_text(slide, title, 0.7, 1.5, 6.6, 2.0,
             size=32, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)

    # Taronga Tracka brand line
    add_text(slide, 'TARONGA TRACKA', 0.7, 3.6, 6, 0.5,
             size=13, bold=True, color=sc, align=PP_ALIGN.LEFT)

    # Outcomes
    oc_text = '  ·  '.join(data['outcomes'])
    add_text(slide, f'NSW Outcomes: {oc_text}', 0.7, 4.1, 8.5, 0.35,
             size=10, bold=False, color=RGBColor(0x88,0xAA,0x99), align=PP_ALIGN.LEFT)

    # Divider line
    add_rect(slide, 0.7, 4.55, 8.6, 0.025, RGBColor(0x2E,0x7D,0x55))

    # Sub info
    timing_desc = 'Prepare students for their Taronga excursion' if timing == 'pre' else 'Consolidate learning after the Taronga excursion'
    add_text(slide, timing_desc, 0.7, 4.7, 8.5, 0.4,
             size=11, bold=False, color=C_MUTED, align=PP_ALIGN.LEFT)

    # Bottom bar
    add_rect(slide, 0, 7.2, 10, 0.3, C_MID)
    add_text(slide, 'taronga.org.au  ·  tarongatracka.com.au', 0.2, 7.22, 9.6, 0.25,
             size=9, bold=False, color=RGBColor(0x88,0xAA,0x99), align=PP_ALIGN.RIGHT)

def make_li_sc_slide(prs, data, timing, subj_color):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_slide_background(slide, prs, C_DEEP)

    # Top bar
    add_rect(slide, 0, 0, 10, 0.08, subj_color)

    # Header
    add_rect(slide, 0, 0.08, 10, 0.9, C_MID)
    add_text(slide, 'LEARNING INTENTION & SUCCESS CRITERIA', 0.3, 0.15, 9.4, 0.7,
             size=14, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)

    # LI box
    add_rect(slide, 0.3, 1.2, 9.4, 1.5, RGBColor(0x0E, 0x2E, 0x1E))
    add_rect(slide, 0.3, 1.2, 0.08, 1.5, subj_color)
    add_text(slide, 'LEARNING INTENTION', 0.5, 1.25, 4, 0.35,
             size=10, bold=True, color=subj_color)
    add_text(slide, data['li'], 0.5, 1.6, 8.9, 0.9,
             size=14, bold=False, color=C_WHITE)

    # SC items
    add_text(slide, 'SUCCESS CRITERIA — I can...', 0.3, 2.95, 9, 0.4,
             size=12, bold=True, color=C_GOLD)
    for i, sc in enumerate(data['sc']):
        y = 3.4 + i * 0.9
        add_rect(slide, 0.3, y, 9.4, 0.75, RGBColor(0x0A, 0x26, 0x1A))
        add_rect(slide, 0.3, y, 0.08, 0.75, subj_color)
        num_box = add_rect(slide, 0.45, y + 0.12, 0.35, 0.35, subj_color)
        add_text(slide, str(i+1), 0.47, y + 0.1, 0.32, 0.4,
                 size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, sc, 0.95, y + 0.1, 8.5, 0.55, size=12, color=C_WHITE)

    # Bottom bar
    add_rect(slide, 0, 7.2, 10, 0.3, C_MID)

def make_video_slide(prs, subject, timing, subj_color):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_slide_background(slide, prs, RGBColor(0x03, 0x10, 0x08))

    add_rect(slide, 0, 0, 10, 0.08, subj_color)

    # Play button visual
    add_rect(slide, 3.2, 1.0, 3.6, 2.8, RGBColor(0x0A, 0x26, 0x1A))

    # Triangle play icon (approximate with text symbol)
    add_text(slide, '▶', 4.45, 1.5, 1.1, 1.1,
             size=48, bold=False, color=subj_color, align=PP_ALIGN.CENTER)

    add_text(slide, 'VIDEO PLACEHOLDER', 0, 3.95, 10, 0.5,
             size=20, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    if timing == 'pre':
        vid_desc = 'Insert a short clip introducing the animals students will observe today.\nSuggested: Taronga keeper video, National Geographic clip, or animal documentary excerpt.\nAim for 2–4 minutes — pause at key moments to discuss.'
    else:
        vid_desc = 'Insert a short clip reviewing the animals students observed at the zoo.\nSuggested: student video from the visit, Taronga social media, or conservation documentary.\nAim for 2–4 minutes — use as a reflection anchor.'

    add_text(slide, vid_desc, 0.8, 4.55, 8.4, 1.2, size=13, color=C_MUTED,
             align=PP_ALIGN.CENTER)

    add_text(slide, 'Teacher note: embed your chosen video directly into this slide before presenting.',
             0.5, 6.4, 9, 0.4, size=10, color=RGBColor(0x55, 0x88, 0x66), align=PP_ALIGN.CENTER)

    add_rect(slide, 0, 7.2, 10, 0.3, C_MID)

def make_content_slide(prs, heading, bullets, slide_num, total_slides, subj_color, accent=False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_slide_background(slide, prs, C_DEEP)

    add_rect(slide, 0, 0, 10, 0.08, subj_color)

    # Header band
    header_color = subj_color if accent else C_MID
    add_rect(slide, 0, 0.08, 10, 0.85, header_color)
    add_text(slide, heading.upper(), 0.35, 0.18, 9, 0.65,
             size=17, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)

    # Slide counter
    add_text(slide, f'{slide_num} / {total_slides}', 9.2, 0.2, 0.6, 0.35,
             size=9, color=RGBColor(0xA8,0xC4,0xB2), align=PP_ALIGN.RIGHT)

    # Bullet content
    y_start = 1.15
    for i, bullet in enumerate(bullets):
        y = y_start + i * 1.1
        add_rect(slide, 0.3, y, 9.4, 0.95, RGBColor(0x0A, 0x26, 0x1A))
        add_rect(slide, 0.3, y, 0.06, 0.95, subj_color)
        add_text(slide, bullet, 0.5, y + 0.12, 9.0, 0.72, size=14, color=C_WHITE)

    add_rect(slide, 0, 7.2, 10, 0.3, C_MID)

def make_brain_break_slide(prs, stage, subj_color):
    import random
    bb = BRAIN_BREAKS[(stage - 2) % len(BRAIN_BREAKS)]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_slide_background(slide, prs, RGBColor(0x05, 0x1A, 0x10))

    add_rect(slide, 0, 0, 10, 0.08, C_GOLD)

    # BRAIN BREAK header
    add_rect(slide, 0, 0.08, 10, 1.1, C_GOLD)
    add_text(slide, '🧠  BRAIN BREAK', 0.35, 0.15, 9, 0.85,
             size=28, bold=True, color=C_DEEP, align=PP_ALIGN.LEFT)

    # Activity name
    add_text(slide, bb[0].upper(), 0.3, 1.45, 9.4, 0.6,
             size=22, bold=True, color=C_GOLD, align=PP_ALIGN.LEFT)

    # Instruction box
    add_rect(slide, 0.3, 2.2, 9.4, 2.8, RGBColor(0x0A, 0x26, 0x1A))
    add_rect(slide, 0.3, 2.2, 0.1, 2.8, C_GOLD)
    add_text(slide, bb[1], 0.55, 2.35, 8.9, 2.5, size=15, color=C_WHITE)

    add_text(slide, '⏱  Allow 3–5 minutes  ·  Everyone participates!', 0.3, 5.2, 9.4, 0.5,
             size=12, bold=True, color=RGBColor(0x88,0xAA,0x77), align=PP_ALIGN.CENTER)

    add_rect(slide, 0, 7.2, 10, 0.3, C_MID)

def make_applied_learning_slide(prs, subject, stage, timing, data, subj_color):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_slide_background(slide, prs, C_DEEP)

    add_rect(slide, 0, 0, 10, 0.08, subj_color)
    add_rect(slide, 0, 0.08, 10, 1.0, C_MID)

    if timing == 'pre':
        add_text(slide, 'PLANNING SESSION — Getting Ready for the Zoo', 0.35, 0.2, 9.3, 0.75,
                 size=16, bold=True, color=C_WHITE)
    else:
        add_text(slide, 'ACTION PROJECT — Apply Your Learning', 0.35, 0.2, 9.3, 0.75,
                 size=16, bold=True, color=C_WHITE)

    if timing == 'pre':
        # Planning mission
        add_text(slide, 'Your Mission at Taronga Zoo', 0.3, 1.3, 9.4, 0.45,
                 size=16, bold=True, color=C_GOLD)
        add_text(slide, data['tracka_focus'], 0.3, 1.8, 9.4, 0.6, size=13, color=C_WHITE)

        add_text(slide, 'Planning Checklist', 0.3, 2.6, 9.4, 0.4,
                 size=15, bold=True, color=subj_color)
        checks = [
            'Choose the animals you most want to observe — write them down now',
            'What specific thing will you look for at each animal? (behaviour / feature / data)',
            f'Which prior knowledge from {SUBJ_LABELS[subject]} can you connect to each animal?',
            'How will you earn maximum points? (read the scoring guide on the next slide)',
        ]
        for i, c in enumerate(checks):
            y = 3.1 + i * 0.8
            add_rect(slide, 0.3, y, 0.5, 0.55, RGBColor(0x0A, 0x26, 0x1A))
            add_text(slide, '☐', 0.33, y + 0.05, 0.45, 0.45, size=16, color=subj_color, align=PP_ALIGN.CENTER)
            add_text(slide, c, 0.9, y + 0.07, 8.8, 0.55, size=12, color=C_WHITE)
    else:
        # Action project
        title, desc = data['action']
        add_text(slide, title, 0.3, 1.3, 9.4, 0.45,
                 size=17, bold=True, color=C_GOLD)
        add_rect(slide, 0.3, 1.85, 9.4, 3.0, RGBColor(0x0A, 0x26, 0x1A))
        add_rect(slide, 0.3, 1.85, 0.08, 3.0, subj_color)
        add_text(slide, desc, 0.5, 1.95, 8.9, 2.8, size=12, color=C_WHITE)

        add_text(slide, 'Project formats: Poster · Video · Model · Report · Presentation · Campaign',
                 0.3, 5.05, 9.4, 0.4, size=11, bold=True, color=RGBColor(0x88,0xAA,0x77), align=PP_ALIGN.CENTER)

    add_rect(slide, 0, 7.2, 10, 0.3, C_MID)

def make_exit_ticket_slide(prs, subject, stage, timing, data, subj_color):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_slide_background(slide, prs, C_DEEP)

    add_rect(slide, 0, 0, 10, 0.08, C_GOLD)
    add_rect(slide, 0, 0.08, 10, 1.0, RGBColor(0x1A, 0x14, 0x04))
    add_text(slide, '🎫  EXIT TICKET', 0.35, 0.2, 9, 0.75,
             size=24, bold=True, color=C_GOLD, align=PP_ALIGN.LEFT)

    for i, q in enumerate(data['exit_q']):
        y = 1.35 + i * 2.1
        add_rect(slide, 0.3, y, 9.4, 1.85, RGBColor(0x0A, 0x26, 0x1A))
        add_rect(slide, 0.3, y, 0.08, 1.85, C_GOLD)
        add_text(slide, f'Q{i+1}', 0.48, y + 0.15, 0.5, 0.45,
                 size=16, bold=True, color=C_GOLD)
        add_text(slide, q, 0.5, y + 0.2, 8.9, 1.5, size=13, color=C_WHITE)

    # Before/after indicator
    if timing == 'pre':
        note = 'Hand in before leaving for the zoo — or submit via your teacher\'s instructions.'
    else:
        note = 'Hand in before leaving class — or submit digitally as directed by your teacher.'
    add_text(slide, note, 0.3, 6.6, 9.4, 0.4,
             size=11, color=RGBColor(0x88,0xAA,0x77), align=PP_ALIGN.CENTER)

    add_rect(slide, 0, 7.2, 10, 0.3, C_MID)
    add_text(slide, 'taronga.org.au  ·  TARONGA TRACKA', 0.2, 7.22, 9.6, 0.25,
             size=9, bold=False, color=RGBColor(0x55, 0x88, 0x66), align=PP_ALIGN.RIGHT)

def generate_deck(subject, stage, timing):
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.5)

    data = CONTENT[subject][stage]
    sc = SUBJ_COLORS[subject]

    if timing == 'pre':
        content_slides = data['pre_content']
    else:
        content_slides = data['post_content']

    total = 2 + 1 + len(content_slides) + 1 + 1 + 1  # title + LI + video + content + brain + applied + exit

    # 1 — Title
    make_title_slide(prs, subject, stage, timing, data)

    # 2 — LI & SC
    make_li_sc_slide(prs, data, timing, sc)

    # 3 — Video placeholder
    make_video_slide(prs, subject, timing, sc)

    # 4..N — Content slides
    for i, (heading, bullets) in enumerate(content_slides):
        make_content_slide(prs, heading, bullets, i+4, total, sc, accent=(i==0))

    # Brain break
    make_brain_break_slide(prs, stage, sc)

    # Applied learning
    make_applied_learning_slide(prs, subject, stage, timing, data, sc)

    # Exit ticket
    make_exit_ticket_slide(prs, subject, stage, timing, data, sc)

    filename = f'{timing}-{subject}-s{stage}.pptx'
    out_path = os.path.join(OUT_DIR, filename)
    prs.save(out_path)
    print(f'  ✓  {filename}')

if __name__ == '__main__':
    print('Generating Taronga Tracka lesson decks...')
    subjects = ['science', 'maths', 'english', 'pdhpe']
    stages   = [2, 3, 4, 5]
    timings  = ['pre', 'post']

    count = 0
    for timing in timings:
        for subject in subjects:
            for stage in stages:
                generate_deck(subject, stage, timing)
                count += 1

    print(f'\nDone — {count} decks written to public/resources/pptx/')
